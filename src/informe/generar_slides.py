"""
Entregable 3 — Diapositivas del pitch (PowerPoint, 10-15 min).

Genera informe/Pitch_Defensa.pptx siguiendo el guion de GUION_PITCH.md, con la
paleta corporativa del proyecto y las cifras leídas de los JSON del pipeline
(si se reejecuta el pipeline, las slides se regeneran consistentes).
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

RUTA_PROYECTO = Path(__file__).resolve().parents[2]
RUTA_REPORTES = RUTA_PROYECTO / "reports"
RUTA_FIGURAS = RUTA_REPORTES / "figures"
RUTA_SALIDA = RUTA_PROYECTO / "informe" / "Pitch_Defensa.pptx"

VERDE_OSCURO = RGBColor(0x1B, 0x43, 0x32)
VERDE = RGBColor(0x2E, 0x7D, 0x32)
VERDE_CLARO = RGBColor(0x66, 0xBB, 0x6A)
AMBAR = RGBColor(0xF9, 0xA8, 0x25)
ROJO = RGBColor(0xC6, 0x28, 0x28)
CREMA = RGBColor(0xF6, 0xF8, 0xF4)
TINTA = RGBColor(0x1C, 0x2B, 0x21)
TINTA_SUAVE = RGBColor(0x5C, 0x6F, 0x62)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

ANCHO, ALTO = Inches(13.333), Inches(7.5)

with open(RUTA_REPORTES / "eda_resultados.json", encoding="utf-8") as f:
    EDA = json.load(f)
with open(RUTA_REPORTES / "modelo_resultados.json", encoding="utf-8") as f:
    MODELO = json.load(f)
RECS = pd.read_csv(RUTA_REPORTES / "recomendaciones_compra.csv")

KPIS = EDA["kpis"]
RF = MODELO["metricas_por_modelo"]["Random Forest"]
ANIO_OBJ = int(RECS["anio_objetivo"].iloc[0])

prs = Presentation()
prs.slide_width, prs.slide_height = ANCHO, ALTO
LAYOUT_VACIO = prs.slide_layouts[6]


def fondo(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def caja_texto(slide, x, y, w, h, texto, tam, color, negrita=False,
               alineacion=PP_ALIGN.LEFT, interlineado=1.0):
    caja = slide.shapes.add_textbox(x, y, w, h)
    tf = caja.text_frame
    tf.word_wrap = True
    lineas = texto.split("\n")
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = linea
        p.alignment = alineacion
        p.line_spacing = interlineado
        for run in p.runs:
            run.font.size = Pt(tam)
            run.font.bold = negrita
            run.font.color.rgb = color
            run.font.name = "Segoe UI"
    return caja


def rectangulo(slide, x, y, w, h, color, radio=False):
    from pptx.enum.shapes import MSO_SHAPE
    forma = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radio else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    forma.fill.solid()
    forma.fill.fore_color.rgb = color
    forma.line.fill.background()
    forma.shadow.inherit = False
    return forma


def barra_titulo(slide, titulo, subtitulo=""):
    rectangulo(slide, 0, 0, ANCHO, Inches(1.15), VERDE_OSCURO)
    caja_texto(slide, Inches(0.55), Inches(0.18), Inches(11.5), Inches(0.6),
               titulo, 26, BLANCO, negrita=True)
    if subtitulo:
        caja_texto(slide, Inches(0.55), Inches(0.72), Inches(11.5), Inches(0.35),
                   subtitulo, 12, RGBColor(0xC8, 0xE6, 0xC9))


def pie(slide, numero):
    caja_texto(slide, Inches(0.55), Inches(7.05), Inches(11), Inches(0.35),
               "AgroComercial del Litoral S.A. · Analítica de Negocios 7A", 9, TINTA_SUAVE)
    caja_texto(slide, Inches(12.3), Inches(7.05), Inches(0.7), Inches(0.35),
               str(numero), 9, TINTA_SUAVE, alineacion=PP_ALIGN.RIGHT)


def tarjeta_kpi(slide, x, y, w, h, valor, etiqueta, color=VERDE):
    tarjeta = rectangulo(slide, x, y, w, h, BLANCO, radio=True)
    tarjeta.line.color.rgb = color
    tarjeta.line.width = Pt(1.5)
    caja_texto(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.6),
               valor, 28, color, negrita=True, alineacion=PP_ALIGN.CENTER)
    caja_texto(slide, x + Inches(0.1), y + Inches(0.78), w - Inches(0.2), h - Inches(0.85),
               etiqueta, 10.5, TINTA_SUAVE, alineacion=PP_ALIGN.CENTER)


def imagen_ajustada(slide, ruta, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(ruta), x, y, width=w, height=h)


# ============================ 1. PORTADA =====================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, VERDE_OSCURO)
rectangulo(s, 0, Inches(6.9), ANCHO, Inches(0.6), VERDE)
caja_texto(s, Inches(1), Inches(1.6), Inches(11.3), Inches(1.6),
           "Dejar de apostar.\nEmpezar a predecir.", 44, BLANCO, negrita=True)
caja_texto(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
           "Sistema de predicción del rendimiento económico de reventa de cosechas", 20,
           RGBColor(0xC8, 0xE6, 0xC9))
caja_texto(s, Inches(1), Inches(4.35), Inches(11.3), Inches(0.5),
           "AgroComercial del Litoral S.A. · Proyecto Integrador — Analítica de Negocios 7A", 14,
           RGBColor(0xA5, 0xD6, 0xA7))
caja_texto(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.9),
           "Integrantes:  _____________  ·  _____________  ·  _____________\n"
           "Docente: Ing. Carlos Manosalvas G.", 13, RGBColor(0xE8, 0xF2, 0xEA))

# ===================== 2. EL PROBLEMA ========================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, "El problema: comprar a ciegas cuesta caro",
             "Hoy las decisiones de compra son reactivas e intuitivas")
caja_texto(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.3),
           "Compramos cosechas apostando a que rendirán como el año pasado. Si la cosecha viene baja, "
           "pagamos caro y no abastecemos (quiebre de stock). Si viene abundante, compramos caro lo que "
           "luego vale menos (sobrestock y merma).", 16, TINTA, interlineado=1.15)
y = Inches(3.2)
tarjeta_kpi(s, Inches(0.8), y, Inches(3.6), Inches(1.7),
            f"{KPIS['KPI_2_variabilidad_rendimiento_pct']['valor']} %",
            "La oferta cambia hasta un tercio\nde un año a otro (CV interanual)", ROJO)
tarjeta_kpi(s, Inches(4.9), y, Inches(3.6), Inches(1.7),
            f"{KPIS['KPI_4_merma_promedio_pct']['valor']} %",
            "Merma promedio por orden:\ndinero que se pudre en bodega", AMBAR)
tarjeta_kpi(s, Inches(9.0), y, Inches(3.6), Inches(1.7),
            f"{KPIS['KPI_5_concentracion_top3_pct']['valor']} %",
            "Del volumen concentrado en solo\n3 cultivos (riesgo de portafolio)", TINTA_SUAVE)
caja_texto(s, Inches(0.8), Inches(5.4), Inches(11.8), Inches(0.9),
           "El problema no es falta de datos —FAO y Banco Mundial publican todo—,\n"
           "es que nadie los estaba usando para decidir.", 17, VERDE, negrita=True,
           alineacion=PP_ALIGN.CENTER)
pie(s, 2)

# ===================== 3. LA SOLUCIÓN ========================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, "La solución: analítica de punta a punta",
             "De los datos oficiales a la decisión de compra")
pasos = [
    ("1 · DATOS", "FAOSTAT (FAO) +\nAPI clima Banco Mundial\n34,113 cosechas · 183 países\n1990–2023", VERDE_OSCURO),
    ("2 · INGENIERÍA", "Pipeline ETL automatizado\nEsquema estrella\n(2 hechos, 4 dimensiones)\nSin datos personales", VERDE),
    ("3 · MODELO", f"Random Forest\nR² = {RF['r2']:.2f}\nError ±{RF['mae_ton_ha']} ton/ha\nValidación cruzada 5-fold", VERDE_CLARO),
    ("4 · DECISIÓN", f"Regla prescriptiva:\naumentar / mantener / reducir\nla compra por cultivo\nPlan {ANIO_OBJ} listo", AMBAR),
]
x = Inches(0.55)
for titulo, cuerpo, color in pasos:
    rectangulo(s, x, Inches(1.7), Inches(2.95), Inches(3.4), color, radio=True)
    caja_texto(s, x + Inches(0.2), Inches(1.95), Inches(2.55), Inches(0.5),
               titulo, 15, BLANCO, negrita=True)
    caja_texto(s, x + Inches(0.2), Inches(2.55), Inches(2.55), Inches(2.4),
               cuerpo, 12.5, BLANCO, interlineado=1.2)
    x += Inches(3.15)
caja_texto(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.1),
           "Entregables: dashboard estratégico (Power BI + versión interactiva), modelo predictivo "
           "funcional, informe técnico de 4 fases y este plan de compras.", 14, TINTA)
pie(s, 3)

# ===================== 4. DATOS Y GOBERNANZA =================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, "Datos oficiales y ética desde el diseño",
             "Fase 2 — Ingeniería de datos y cumplimiento normativo")
vinetas = [
    ("Tres fuentes integradas", "FAOSTAT oficial (rendimiento y pesticidas, hasta 2023-24), API pública de "
     "clima del Banco Mundial (CCKP, consultada en vivo) y el Excel interno de compras."),
    ("Esquema estrella listo para BI", "34,113 observaciones limpias: duplicados, rangos físicos y "
     "consistencia referencial verificados automáticamente en cada ejecución."),
    ("Cero datos personales (LOPDP)", "Proveedores anonimizados por código desde el origen; el pipeline "
     "aborta si detecta columnas con PII."),
    ("Predicción con supervisión humana", "El modelo apoya al comité de compras, no lo reemplaza; toda "
     "predicción se comunica con su margen de error."),
]
y = Inches(1.55)
for titulo, cuerpo in vinetas:
    rectangulo(s, Inches(0.55), y, Inches(0.12), Inches(1.05), VERDE)
    caja_texto(s, Inches(0.9), y - Inches(0.05), Inches(11.8), Inches(0.45),
               titulo, 16, VERDE_OSCURO, negrita=True)
    caja_texto(s, Inches(0.9), y + Inches(0.42), Inches(11.8), Inches(0.7),
               cuerpo, 13, TINTA)
    y += Inches(1.32)
pie(s, 4)

# ===================== 5. HALLAZGOS DESCRIPTIVOS =============================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, "Qué nos dicen 33 años de cosechas",
             "Fase 3 — Analítica descriptiva (Ecuador, 1990-2023)")
imagen_ajustada(s, RUTA_FIGURAS / "02_rendimiento_ecuador.png",
                Inches(0.55), Inches(1.5), w=Inches(12.2))
caja_texto(s, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.7),
           "Papa y plátano lideran el rendimiento; trigo y soya son los más bajos. El rendimiento NO es "
           "estable: por eso el promedio histórico no sirve para planificar compras.", 14, TINTA, negrita=True)
pie(s, 5)

# ===================== 6. DIAGNÓSTICO ========================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, "Correlación no es causalidad",
             "Fase 3 — Analítica diagnóstica: ¿por qué rinde más un cultivo?")
imagen_ajustada(s, RUTA_FIGURAS / "03_correlaciones.png",
                Inches(0.55), Inches(1.5), w=Inches(7.6))
caja_texto(s, Inches(8.5), Inches(1.7), Inches(4.3), Inches(4.6),
           "La lluvia y la temperatura, por sí solas, casi no correlacionan con el rendimiento global.\n\n"
           "¿El clima no importa? Sí importa: su efecto depende del cultivo y del país.\n\n"
           "Los países con más pesticidas rinden más… pero el pesticida delata agricultura tecnificada; "
           "no es la causa por sí solo.\n\n"
           "Conclusión: se necesita un modelo que combine todos los factores a la vez.",
           13.5, TINTA, interlineado=1.15)
pie(s, 6)

# ===================== 7. EL MODELO ==========================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, f"El modelo: Random Forest · R² = {RF['r2']:.2f}",
             "Fase 4 — Analítica predictiva con métricas honestas")
imagen_ajustada(s, RUTA_FIGURAS / "06_comparacion_modelos.png",
                Inches(0.55), Inches(1.55), w=Inches(7.4))
y = Inches(1.7)
for valor, etiqueta in [
    (f"{RF['r2'] * 100:.0f} %", "de la variabilidad del\nrendimiento explicada"),
    (f"±{RF['mae_ton_ha']}", "ton/ha de error típico\n(MAE en datos nunca vistos)"),
    (f"{MODELO['error_por_grupo']['rmse_ecuador']}", "ton/ha RMSE en Ecuador\n(mejor que el promedio global)"),
]:
    tarjeta_kpi(s, Inches(8.5), y, Inches(4.2), Inches(1.5), valor, etiqueta)
    y += Inches(1.7)
pie(s, 7)

# ===================== 8. XAI ================================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, "No es una caja negra",
             "IA Explicable (XAI): qué variables pesan en la predicción")
imagen_ajustada(s, RUTA_FIGURAS / "08_importancia_variables.png",
                Inches(0.55), Inches(1.6), w=Inches(8.2))
caja_texto(s, Inches(9.1), Inches(1.9), Inches(3.7), Inches(4.3),
           "El modelo decide como un agrónomo:\n\n"
           "1º ubica el cultivo\n\n2º considera el país (suelo, tecnología)\n\n"
           "3º ajusta por clima y manejo del año\n\n"
           "…pero con la memoria de 34,113 cosechas.",
           14, TINTA, interlineado=1.2)
pie(s, 8)

# ===================== 9. PLAN DE COMPRAS ====================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, f"El entregable: plan de compras {ANIO_OBJ}",
             "Pilar prescriptivo — regla de decisión sobre el índice de oferta")
colores_dec = {"AUMENTAR compra": VERDE, "MANTENER volumen": AMBAR, "REDUCIR compra": ROJO}
etiquetas_dec = {"AUMENTAR compra": "▲ +15 %", "MANTENER volumen": "= mantener", "REDUCIR compra": "▼ −15 %"}
y = Inches(1.5)
alto_fila = Inches(0.56)
caja_texto(s, Inches(0.8), y, Inches(3.2), Inches(0.4), "Cultivo", 13, TINTA_SUAVE, negrita=True)
caja_texto(s, Inches(4.2), y, Inches(2.4), Inches(0.4), "Predicho (ton/ha)", 13, TINTA_SUAVE, negrita=True)
caja_texto(s, Inches(6.8), y, Inches(2.0), Inches(0.4), "Índice oferta", 13, TINTA_SUAVE, negrita=True)
caja_texto(s, Inches(9.0), y, Inches(3.5), Inches(0.4), "Decisión", 13, TINTA_SUAVE, negrita=True)
y += Inches(0.5)
for _, fila in RECS.iterrows():
    color = colores_dec[fila["decision"]]
    rectangulo(s, Inches(0.55), y, Inches(0.1), alto_fila - Inches(0.08), color)
    caja_texto(s, Inches(0.8), y, Inches(3.2), alto_fila, fila["cultivo"], 14, TINTA, negrita=True)
    caja_texto(s, Inches(4.2), y, Inches(2.4), alto_fila, f"{fila['rendimiento_predicho_ton_ha']}", 14, TINTA)
    caja_texto(s, Inches(6.8), y, Inches(2.0), alto_fila, f"{fila['indice_oferta']}", 14, TINTA)
    caja_texto(s, Inches(9.0), y, Inches(3.5), alto_fila, etiquetas_dec[fila["decision"]], 14, color, negrita=True)
    y += alto_fila
pie(s, 9)

# ===================== 10. LÍMITES Y ÉTICA ===================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, CREMA)
barra_titulo(s, "Honestidad que da credibilidad",
             "Sesgos, límites y cómo los mitigamos")
limites = [
    ("Sesgo de representación", "Los cultivos con menos historia (plátano, papa en ciertos países) tienen "
     "mayor error; sus recomendaciones se entregan con advertencia explícita de incertidumbre."),
    ("Rezago del dato oficial", "FAOSTAT publica con 1-2 años de rezago (histórico hasta 2023). El pipeline "
     "reentrena con un comando cuando la FAO actualiza."),
    ("Compras internas simuladas", "El rendimiento y el clima son reales; las compras son un escenario "
     "académico declarado como tal."),
    ("Supervisión humana", "La decisión final es del comité: el modelo aporta evidencia, no órdenes."),
]
y = Inches(1.55)
for titulo, cuerpo in limites:
    rectangulo(s, Inches(0.55), y, Inches(0.12), Inches(1.05), AMBAR)
    caja_texto(s, Inches(0.9), y - Inches(0.05), Inches(11.8), Inches(0.45),
               titulo, 16, VERDE_OSCURO, negrita=True)
    caja_texto(s, Inches(0.9), y + Inches(0.42), Inches(11.8), Inches(0.7),
               cuerpo, 13, TINTA)
    y += Inches(1.32)
pie(s, 10)

# ===================== 11. CIERRE ============================================
s = prs.slides.add_slide(LAYOUT_VACIO)
fondo(s, VERDE_OSCURO)
caja_texto(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.8),
           "Hoy AgroComercial decide con evidencia:\nqué comprar, cuánto y con qué confianza.",
           34, BLANCO, negrita=True)
tarjetas_cierre = [
    (f"{KPIS['KPI_3_margen_bruto_pct']['valor']} %", "de margen bruto,\ndefendido con datos"),
    (f"R² {RF['r2']:.2f}", "modelo validado en\ndatos nunca vistos"),
    ("34,113", "cosechas de respaldo para\nla intuición del comprador"),
]
for posicion, (valor, etiqueta) in enumerate(tarjetas_cierre):
    tarjeta_kpi(s, Inches(1.0) + Emu(int(Inches(4.0)) * posicion), Inches(4.0),
                Inches(3.6), Inches(1.7), valor, etiqueta)
caja_texto(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.6),
           "Gracias — ¿preguntas?", 20, RGBColor(0xC8, 0xE6, 0xC9), alineacion=PP_ALIGN.CENTER)

prs.save(RUTA_SALIDA)
print(f"OK — {len(prs.slides._sldIdLst)} diapositivas en {RUTA_SALIDA}")
